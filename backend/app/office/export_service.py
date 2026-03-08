"""Import/export service for office documents (PPTX, XLSX, DOCX, PDF)."""

import io
import json
import logging
from typing import Any

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# PPTX Export
# ---------------------------------------------------------------------------

def export_pptx(title: str, content_json: dict | None) -> bytes:
    """Export presentation content_json to PPTX bytes."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = []
    if content_json and "slides" in content_json:
        slides_data = content_json["slides"]

    if not slides_data:
        # Add a blank title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        if slide.placeholders[0]:
            slide.placeholders[0].text = title
    else:
        for slide_data in slides_data:
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Extract text from Tiptap JSON content
            tiptap_content = slide_data.get("content")
            if tiptap_content and isinstance(tiptap_content, dict):
                text_runs = _extract_tiptap_text_runs(tiptap_content)
                if text_runs:
                    txBox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.5),
                        Inches(12.333), Inches(6.5)
                    )
                    tf = txBox.text_frame
                    tf.word_wrap = True

                    first = True
                    for run_info in text_runs:
                        if first:
                            p = tf.paragraphs[0]
                            first = False
                        else:
                            if run_info.get("new_paragraph", False):
                                p = tf.add_paragraph()

                        run = p.add_run()
                        run.text = run_info.get("text", "")

                        font_size = run_info.get("font_size")
                        if font_size:
                            run.font.size = Pt(font_size)
                        if run_info.get("bold"):
                            run.font.bold = True
                        if run_info.get("italic"):
                            run.font.italic = True
                        if run_info.get("underline"):
                            run.font.underline = True

                        alignment = run_info.get("alignment")
                        if alignment == "center":
                            p.alignment = PP_ALIGN.CENTER
                        elif alignment == "right":
                            p.alignment = PP_ALIGN.RIGHT

            # Add speaker notes
            notes = slide_data.get("speakerNotes", "")
            if notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _extract_tiptap_text_runs(node: dict, depth: int = 0) -> list[dict]:
    """Recursively extract text runs from Tiptap/ProseMirror JSON."""
    runs = []
    node_type = node.get("type", "")

    if node_type == "text":
        text = node.get("text", "")
        marks = node.get("marks", [])
        run: dict[str, Any] = {"text": text}
        for mark in marks:
            mark_type = mark.get("type", "")
            if mark_type == "bold":
                run["bold"] = True
            elif mark_type == "italic":
                run["italic"] = True
            elif mark_type == "underline":
                run["underline"] = True
        return [run]

    attrs = node.get("attrs", {})
    content = node.get("content", [])

    is_heading = node_type == "heading"
    level = attrs.get("level", 2) if is_heading else None
    alignment = attrs.get("textAlign")
    is_block = node_type in ("paragraph", "heading", "bulletList", "orderedList", "listItem", "blockquote")

    child_runs = []
    for child in content:
        child_runs.extend(_extract_tiptap_text_runs(child, depth + 1))

    if is_block and child_runs:
        child_runs[0]["new_paragraph"] = True
        if is_heading:
            font_sizes = {1: 36, 2: 28, 3: 22}
            for r in child_runs:
                r["font_size"] = font_sizes.get(level, 22)
                r["bold"] = True
        if alignment:
            for r in child_runs:
                r["alignment"] = alignment

    runs.extend(child_runs)

    if node_type in ("paragraph", "heading") and not child_runs:
        runs.append({"text": "", "new_paragraph": True})

    return runs


# ---------------------------------------------------------------------------
# PPTX Import
# ---------------------------------------------------------------------------

def import_pptx(file_bytes: bytes) -> dict:
    """Import PPTX bytes to presentation content_json."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []

    for slide in prs.slides:
        tiptap_content: dict[str, Any] = {"type": "doc", "content": []}

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_node: dict[str, Any] = {"type": "paragraph", "content": []}

                    for run in paragraph.runs:
                        text_node: dict[str, Any] = {"type": "text", "text": run.text}
                        marks = []
                        if run.font.bold:
                            marks.append({"type": "bold"})
                        if run.font.italic:
                            marks.append({"type": "italic"})
                        if run.font.underline:
                            marks.append({"type": "underline"})
                        if marks:
                            text_node["marks"] = marks
                        if run.text:
                            para_node["content"].append(text_node)

                    # Check if this looks like a heading (large font)
                    is_heading = False
                    for run in paragraph.runs:
                        if run.font.size and run.font.size.pt >= 24:
                            is_heading = True
                            break

                    if is_heading:
                        para_node["type"] = "heading"
                        para_node["attrs"] = {"level": 1}

                    if para_node["content"]:
                        tiptap_content["content"].append(para_node)

        # Get speaker notes
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
        except Exception:
            pass

        slides.append({
            "content": tiptap_content if tiptap_content["content"] else None,
            "speakerNotes": notes,
            "layout": "blank",
            "transition": "none",
            "backgroundColor": "#ffffff",
        })

    if not slides:
        slides = [{"content": None, "speakerNotes": "", "layout": "blank", "transition": "none", "backgroundColor": "#ffffff"}]

    return {"slides": slides}


# ---------------------------------------------------------------------------
# XLSX Export
# ---------------------------------------------------------------------------

def export_xlsx(title: str, content_json: dict | None) -> bytes:
    """Export spreadsheet content_json to XLSX bytes."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    sheets_data = []
    if content_json and "sheets" in content_json:
        sheets_data = content_json["sheets"]

    if not sheets_data:
        ws = wb.active
        ws.title = title or "Sheet 1"
    else:
        for i, sheet_data in enumerate(sheets_data):
            if i == 0:
                ws = wb.active
            else:
                ws = wb.create_sheet()

            ws.title = sheet_data.get("name", f"Sheet {i + 1}")
            cells = sheet_data.get("cells", {})

            for cell_ref, cell_data in cells.items():
                if not isinstance(cell_data, dict):
                    continue

                try:
                    cell = ws[cell_ref]
                except Exception:
                    continue

                value = cell_data.get("value", "")
                formula = cell_data.get("formula", "")

                if formula and formula.startswith("="):
                    cell.value = formula
                elif value is not None:
                    # Try numeric
                    if isinstance(value, (int, float)):
                        cell.value = value
                    else:
                        try:
                            cell.value = float(value)
                        except (ValueError, TypeError):
                            cell.value = str(value) if value else ""

                # Apply formatting
                fmt = cell_data.get("format", {})
                if isinstance(fmt, dict):
                    font_kwargs = {}
                    if fmt.get("bold"):
                        font_kwargs["bold"] = True
                    if fmt.get("italic"):
                        font_kwargs["italic"] = True
                    if fmt.get("underline"):
                        font_kwargs["underline"] = "single"
                    if fmt.get("fontSize"):
                        font_kwargs["size"] = fmt["fontSize"]
                    if fmt.get("color"):
                        color = fmt["color"].lstrip("#")
                        font_kwargs["color"] = color
                    if font_kwargs:
                        cell.font = Font(**font_kwargs)

                    bg = fmt.get("backgroundColor")
                    if bg:
                        bg_color = bg.lstrip("#")
                        cell.fill = PatternFill(start_color=bg_color, end_color=bg_color, fill_type="solid")

                    align = fmt.get("textAlign")
                    if align:
                        cell.alignment = Alignment(horizontal=align)

                # Column width heuristic
                col_letter = cell_ref[0] if cell_ref[0].isalpha() else get_column_letter(cell.column)
                current_width = ws.column_dimensions[col_letter].width or 8
                new_width = max(current_width, min(len(str(cell.value or "")) + 2, 50))
                ws.column_dimensions[col_letter].width = new_width

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# XLSX Import
# ---------------------------------------------------------------------------

def import_xlsx(file_bytes: bytes) -> dict:
    """Import XLSX bytes to spreadsheet content_json."""
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    wb = load_workbook(io.BytesIO(file_bytes), data_only=False)
    sheets = []

    for ws in wb.worksheets:
        cells = {}

        for row in ws.iter_rows():
            for cell in row:
                if cell.value is None:
                    continue

                col_letter = get_column_letter(cell.column)
                cell_ref = f"{col_letter}{cell.row}"

                cell_data: dict[str, Any] = {}

                # Value and formula
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    cell_data["formula"] = cell.value
                    cell_data["value"] = ""
                else:
                    cell_data["value"] = cell.value if cell.value is not None else ""

                # Format
                fmt: dict[str, Any] = {}
                if cell.font:
                    if cell.font.bold:
                        fmt["bold"] = True
                    if cell.font.italic:
                        fmt["italic"] = True
                    if cell.font.underline:
                        fmt["underline"] = True
                    if cell.font.size:
                        fmt["fontSize"] = cell.font.size
                    if cell.font.color and cell.font.color.rgb and cell.font.color.rgb != "00000000":
                        fmt["color"] = f"#{cell.font.color.rgb[-6:]}"

                if cell.fill and cell.fill.start_color and cell.fill.start_color.rgb:
                    rgb = cell.fill.start_color.rgb
                    if isinstance(rgb, str) and rgb != "00000000":
                        fmt["backgroundColor"] = f"#{rgb[-6:]}"

                if cell.alignment and cell.alignment.horizontal:
                    fmt["textAlign"] = cell.alignment.horizontal

                if fmt:
                    cell_data["format"] = fmt

                cells[cell_ref] = cell_data

        sheets.append({
            "name": ws.title,
            "cells": cells,
            "columnWidths": {},
            "rowHeights": {},
        })

    if not sheets:
        sheets = [{"name": "Sheet 1", "cells": {}, "columnWidths": {}, "rowHeights": {}}]

    return {"sheets": sheets, "activeSheet": 0}


# ---------------------------------------------------------------------------
# DOCX Export
# ---------------------------------------------------------------------------

def export_docx(title: str, content_json: dict | None) -> bytes:
    """Export document content_json (Tiptap JSON) to DOCX bytes."""
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Set default font
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Calibri"
    font.size = Pt(11)

    if not content_json:
        doc.add_heading(title, level=0)
    else:
        content_nodes = content_json.get("content", [])
        _render_tiptap_to_docx(doc, content_nodes)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _render_tiptap_to_docx(doc: Any, nodes: list[dict]) -> None:
    """Render Tiptap content nodes into a python-docx Document."""
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    for node in nodes:
        node_type = node.get("type", "")
        attrs = node.get("attrs", {})
        content = node.get("content", [])

        if node_type == "heading":
            level = attrs.get("level", 1)
            text = _tiptap_plain_text(content)
            doc.add_heading(text, level=min(level, 4))

        elif node_type == "paragraph":
            p = doc.add_paragraph()
            _add_tiptap_runs(p, content)
            align = attrs.get("textAlign")
            if align == "center":
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            elif align == "right":
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT

        elif node_type == "bulletList":
            for item in content:
                if item.get("type") == "listItem":
                    item_content = item.get("content", [])
                    for child in item_content:
                        if child.get("type") == "paragraph":
                            p = doc.add_paragraph(style="List Bullet")
                            _add_tiptap_runs(p, child.get("content", []))

        elif node_type == "orderedList":
            for item in content:
                if item.get("type") == "listItem":
                    item_content = item.get("content", [])
                    for child in item_content:
                        if child.get("type") == "paragraph":
                            p = doc.add_paragraph(style="List Number")
                            _add_tiptap_runs(p, child.get("content", []))

        elif node_type == "blockquote":
            for child in content:
                if child.get("type") == "paragraph":
                    p = doc.add_paragraph(style="Quote")
                    _add_tiptap_runs(p, child.get("content", []))

        elif node_type == "horizontalRule":
            p = doc.add_paragraph()
            p.add_run("_" * 60)

        elif node_type == "table":
            rows_data = [r for r in content if r.get("type") == "tableRow"]
            if rows_data:
                num_cols = max(len(r.get("content", [])) for r in rows_data)
                table = doc.add_table(rows=len(rows_data), cols=num_cols)
                table.style = "Table Grid"

                for row_idx, row_node in enumerate(rows_data):
                    cells = row_node.get("content", [])
                    for col_idx, cell_node in enumerate(cells):
                        if col_idx < num_cols:
                            cell = table.cell(row_idx, col_idx)
                            cell_content = cell_node.get("content", [])
                            cell_text = _tiptap_plain_text_from_nodes(cell_content)
                            cell.text = cell_text


def _add_tiptap_runs(paragraph: Any, nodes: list[dict]) -> None:
    """Add formatted runs to a python-docx paragraph from Tiptap inline nodes."""
    from docx.shared import Pt

    for node in nodes:
        if node.get("type") == "text":
            text = node.get("text", "")
            run = paragraph.add_run(text)
            marks = node.get("marks", [])
            for mark in marks:
                mark_type = mark.get("type", "")
                if mark_type == "bold":
                    run.bold = True
                elif mark_type == "italic":
                    run.italic = True
                elif mark_type == "underline":
                    run.underline = True
        elif node.get("type") == "hardBreak":
            paragraph.add_run().add_break()


def _tiptap_plain_text(nodes: list[dict]) -> str:
    """Extract plain text from Tiptap inline content nodes."""
    parts = []
    for node in nodes:
        if node.get("type") == "text":
            parts.append(node.get("text", ""))
    return "".join(parts)


def _tiptap_plain_text_from_nodes(nodes: list[dict]) -> str:
    """Extract plain text from block-level Tiptap nodes."""
    parts = []
    for node in nodes:
        if node.get("type") == "text":
            parts.append(node.get("text", ""))
        elif node.get("type") == "paragraph":
            parts.append(_tiptap_plain_text(node.get("content", [])))
    return " ".join(parts)


# ---------------------------------------------------------------------------
# DOCX Import
# ---------------------------------------------------------------------------

def import_docx(file_bytes: bytes) -> dict:
    """Import DOCX bytes to document content_json (Tiptap JSON)."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    tiptap_nodes: list[dict] = []

    for para in doc.paragraphs:
        style_name = (para.style.name or "").lower()

        if "heading" in style_name:
            # Extract heading level
            level = 1
            for ch in style_name:
                if ch.isdigit():
                    level = int(ch)
                    break

            heading_node: dict[str, Any] = {
                "type": "heading",
                "attrs": {"level": min(level, 3)},
                "content": _docx_runs_to_tiptap(para.runs),
            }
            if heading_node["content"]:
                tiptap_nodes.append(heading_node)

        elif "list bullet" in style_name or "list number" in style_name:
            # Wrap in bullet/ordered list
            list_type = "bulletList" if "bullet" in style_name else "orderedList"
            item_node = {
                "type": list_type,
                "content": [{
                    "type": "listItem",
                    "content": [{
                        "type": "paragraph",
                        "content": _docx_runs_to_tiptap(para.runs),
                    }],
                }],
            }
            if item_node["content"][0]["content"][0]["content"]:
                tiptap_nodes.append(item_node)

        else:
            para_node: dict[str, Any] = {
                "type": "paragraph",
                "content": _docx_runs_to_tiptap(para.runs),
            }
            # Only add if has content or is intentionally empty
            tiptap_nodes.append(para_node)

    # Handle tables
    for table in doc.tables:
        table_node: dict[str, Any] = {
            "type": "table",
            "content": [],
        }
        for row in table.rows:
            row_node: dict[str, Any] = {
                "type": "tableRow",
                "content": [],
            }
            for cell in row.cells:
                cell_text = cell.text
                cell_node: dict[str, Any] = {
                    "type": "tableCell",
                    "content": [{
                        "type": "paragraph",
                        "content": [{"type": "text", "text": cell_text}] if cell_text else [],
                    }],
                }
                row_node["content"].append(cell_node)
            table_node["content"].append(row_node)
        tiptap_nodes.append(table_node)

    return {"type": "doc", "content": tiptap_nodes}


def _docx_runs_to_tiptap(runs: list) -> list[dict]:
    """Convert python-docx runs to Tiptap text nodes with marks."""
    nodes = []
    for run in runs:
        if not run.text:
            continue
        text_node: dict[str, Any] = {"type": "text", "text": run.text}
        marks = []
        if run.bold:
            marks.append({"type": "bold"})
        if run.italic:
            marks.append({"type": "italic"})
        if run.underline:
            marks.append({"type": "underline"})
        if marks:
            text_node["marks"] = marks
        nodes.append(text_node)
    return nodes


# ---------------------------------------------------------------------------
# PDF Export (simple HTML-based approach)
# ---------------------------------------------------------------------------

def export_pdf_html(title: str, content_json: dict | None, doc_type: str = "document") -> str:
    """Generate an HTML string suitable for browser-side PDF printing.

    Returns HTML that can be rendered and printed to PDF on the frontend.
    For server-side PDF, we would use reportlab, but browser print is simpler.
    """
    if doc_type == "spreadsheet":
        return _spreadsheet_to_html(title, content_json)
    elif doc_type == "presentation":
        return _presentation_to_html(title, content_json)
    else:
        return _document_to_html(title, content_json)


def _document_to_html(title: str, content_json: dict | None) -> str:
    """Convert Tiptap document JSON to HTML."""
    html_parts = [
        '<!DOCTYPE html><html><head><meta charset="utf-8">',
        f'<title>{title}</title>',
        '<style>body{font-family:Calibri,sans-serif;max-width:800px;margin:40px auto;padding:0 20px;line-height:1.6;}',
        'h1{font-size:2em}h2{font-size:1.5em}h3{font-size:1.2em}',
        'table{border-collapse:collapse;width:100%;margin:1em 0}td,th{border:1px solid #ddd;padding:8px}',
        'th{background:#f5f5f5}blockquote{border-left:3px solid #ddd;padding-left:1em;color:#666}',
        '@media print{body{margin:0;padding:20px}}</style></head><body>',
    ]

    if content_json:
        nodes = content_json.get("content", [])
        for node in nodes:
            html_parts.append(_tiptap_node_to_html(node))
    else:
        html_parts.append(f"<h1>{title}</h1>")

    html_parts.append('</body></html>')
    return "".join(html_parts)


def _tiptap_node_to_html(node: dict) -> str:
    """Convert a single Tiptap node to HTML."""
    node_type = node.get("type", "")
    attrs = node.get("attrs", {})
    content = node.get("content", [])

    if node_type == "text":
        text = node.get("text", "")
        marks = node.get("marks", [])
        for mark in marks:
            mt = mark.get("type", "")
            if mt == "bold":
                text = f"<strong>{text}</strong>"
            elif mt == "italic":
                text = f"<em>{text}</em>"
            elif mt == "underline":
                text = f"<u>{text}</u>"
            elif mt == "highlight":
                text = f'<mark>{text}</mark>'
        return text

    inner = "".join(_tiptap_node_to_html(c) for c in content)

    align = attrs.get("textAlign", "")
    style = f' style="text-align:{align}"' if align else ""

    if node_type == "heading":
        level = attrs.get("level", 1)
        return f"<h{level}{style}>{inner}</h{level}>"
    elif node_type == "paragraph":
        return f"<p{style}>{inner}</p>"
    elif node_type == "bulletList":
        return f"<ul>{inner}</ul>"
    elif node_type == "orderedList":
        return f"<ol>{inner}</ol>"
    elif node_type == "listItem":
        return f"<li>{inner}</li>"
    elif node_type == "blockquote":
        return f"<blockquote>{inner}</blockquote>"
    elif node_type == "horizontalRule":
        return "<hr>"
    elif node_type == "table":
        return f"<table>{inner}</table>"
    elif node_type == "tableRow":
        return f"<tr>{inner}</tr>"
    elif node_type == "tableHeader":
        return f"<th>{inner}</th>"
    elif node_type == "tableCell":
        return f"<td>{inner}</td>"
    elif node_type == "image":
        src = attrs.get("src", "")
        return f'<img src="{src}" style="max-width:100%">'
    elif node_type == "hardBreak":
        return "<br>"

    return inner


def _spreadsheet_to_html(title: str, content_json: dict | None) -> str:
    """Convert spreadsheet JSON to HTML table."""
    html_parts = [
        '<!DOCTYPE html><html><head><meta charset="utf-8">',
        f'<title>{title}</title>',
        '<style>body{font-family:Calibri,sans-serif;margin:20px}',
        'table{border-collapse:collapse;width:100%}td,th{border:1px solid #ddd;padding:6px 10px;text-align:left}',
        'th{background:#f0f0f0;font-weight:bold}',
        '@media print{body{margin:0}}</style></head><body>',
        f'<h2>{title}</h2>',
    ]

    if content_json and "sheets" in content_json:
        for sheet in content_json["sheets"]:
            cells = sheet.get("cells", {})
            if not cells:
                continue

            html_parts.append(f'<h3>{sheet.get("name", "Sheet")}</h3><table>')

            # Find dimensions
            max_row = 0
            max_col = 0
            for ref in cells:
                col_str = ""
                row_str = ""
                for ch in ref:
                    if ch.isalpha():
                        col_str += ch
                    else:
                        row_str += ch
                col_num = 0
                for ch in col_str.upper():
                    col_num = col_num * 26 + (ord(ch) - ord('A') + 1)
                row_num = int(row_str) if row_str else 0
                max_row = max(max_row, row_num)
                max_col = max(max_col, col_num)

            max_row = min(max_row, 500)
            max_col = min(max_col, 26)

            for r in range(1, max_row + 1):
                html_parts.append("<tr>")
                for c in range(1, max_col + 1):
                    col_letter = chr(ord('A') + c - 1) if c <= 26 else 'A'
                    ref = f"{col_letter}{r}"
                    cell = cells.get(ref, {})
                    val = cell.get("value", "") if isinstance(cell, dict) else ""
                    tag = "th" if r == 1 else "td"
                    html_parts.append(f"<{tag}>{val}</{tag}>")
                html_parts.append("</tr>")

            html_parts.append("</table>")

    html_parts.append('</body></html>')
    return "".join(html_parts)


def _presentation_to_html(title: str, content_json: dict | None) -> str:
    """Convert presentation JSON to HTML (one page per slide)."""
    html_parts = [
        '<!DOCTYPE html><html><head><meta charset="utf-8">',
        f'<title>{title}</title>',
        '<style>body{font-family:Calibri,sans-serif;margin:0}',
        '.slide{width:100%;min-height:100vh;padding:60px;box-sizing:border-box;page-break-after:always;display:flex;align-items:center;justify-content:center}',
        '.slide-content{max-width:900px;width:100%}',
        'h1{font-size:3em}h2{font-size:2em}',
        '@media print{.slide{min-height:auto;page-break-after:always}}</style></head><body>',
    ]

    if content_json and "slides" in content_json:
        for i, slide_data in enumerate(content_json["slides"]):
            bg = slide_data.get("backgroundColor", "#ffffff")
            html_parts.append(f'<div class="slide" style="background:{bg}">')
            html_parts.append('<div class="slide-content">')

            tiptap = slide_data.get("content")
            if tiptap and isinstance(tiptap, dict):
                nodes = tiptap.get("content", [])
                for node in nodes:
                    html_parts.append(_tiptap_node_to_html(node))

            html_parts.append('</div></div>')

    html_parts.append('</body></html>')
    return "".join(html_parts)


# ---------------------------------------------------------------------------
# CSV Export (for spreadsheets)
# ---------------------------------------------------------------------------

def export_csv(content_json: dict | None, sheet_index: int = 0) -> str:
    """Export a single spreadsheet sheet to CSV string."""
    import csv

    if not content_json or "sheets" not in content_json:
        return ""

    sheets = content_json["sheets"]
    if sheet_index >= len(sheets):
        return ""

    sheet = sheets[sheet_index]
    cells = sheet.get("cells", {})

    if not cells:
        return ""

    # Find dimensions
    max_row = 0
    max_col = 0
    for ref in cells:
        col_str = ""
        row_str = ""
        for ch in ref:
            if ch.isalpha():
                col_str += ch
            else:
                row_str += ch
        col_num = 0
        for ch in col_str.upper():
            col_num = col_num * 26 + (ord(ch) - ord('A') + 1)
        row_num = int(row_str) if row_str else 0
        max_row = max(max_row, row_num)
        max_col = max(max_col, col_num)

    buf = io.StringIO()
    writer = csv.writer(buf)

    for r in range(1, max_row + 1):
        row_data = []
        for c in range(1, max_col + 1):
            col_letter = chr(ord('A') + c - 1) if c <= 26 else 'A'
            ref = f"{col_letter}{r}"
            cell = cells.get(ref, {})
            val = cell.get("value", "") if isinstance(cell, dict) else ""
            row_data.append(str(val) if val is not None else "")
        writer.writerow(row_data)

    return buf.getvalue()
